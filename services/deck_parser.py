"""willab beta — slide-deck parser (UX Wave 4 §S / BE-S2).

Turns an uploaded .pptx/.pdf into:
  (a) per-slide {title, body} TEXT — for the editable form + the spoken-vs-
      slide stickiness/compatibility analysis, and
  (b) a single PDF the FE renders with PDF.js (one render path).

PPTX is converted to PDF once via headless LibreOffice (`soffice`); PDF passes
through. Text is extracted from the source: python-pptx for clean pptx titles,
pypdf for pdf (first line → title, rest → body, best-effort with warnings).

Heavy/optional deps (python-pptx, pypdf) are LAZY-imported inside the
functions so the route layer imports this module without them present.
LibreOffice is a system package (nixpacks aptPkgs); warm it via
gunicorn_conf.py so the first real upload doesn't eat the soffice cold start.
"""
from __future__ import annotations

import logging
import os
import shutil
import subprocess
import tempfile
from io import BytesIO

logger = logging.getLogger(__name__)

_SOFFICE_TIMEOUT_SEC = 90
MAX_SLIDES = 60
_TITLE_CAP = 200
_BODY_CAP = 2000

SUPPORTED_EXTS = (".pptx", ".pdf")


class DeckParseError(Exception):
    """Raised on unsupported type / unparseable file / conversion failure.
    The route maps it to 415/422."""


def _ext(filename: str) -> str:
    return (os.path.splitext(filename or "")[1] or "").lower()


def _clip(s, n: int) -> str:
    return (s or "").strip()[:n]


def extract_deck(file_bytes: bytes, filename: str) -> dict:
    """→ {slides: [{title, body}], source: "pptx"|"pdf", warnings: [str],
         pdf_bytes: bytes}. Raises DeckParseError on unsupported/unparseable."""
    ext = _ext(filename)
    if ext == ".pdf":
        slides, warnings = _extract_pdf_text(file_bytes)
        return {"slides": slides, "source": "pdf", "warnings": warnings,
                "pdf_bytes": file_bytes}
    if ext == ".pptx":
        slides, warnings = _extract_pptx_text(file_bytes)
        pdf_bytes = _pptx_to_pdf(file_bytes)
        return {"slides": slides, "source": "pptx", "warnings": warnings,
                "pdf_bytes": pdf_bytes}
    raise DeckParseError(f"unsupported file type: {ext or 'unknown'}")


def _extract_pptx_text(file_bytes: bytes):
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise DeckParseError(f"pptx parser unavailable: {e}")
    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as e:
        raise DeckParseError(f"could not open pptx: {e}")
    slides: list[dict] = []
    warnings: list[str] = []
    for i, slide in enumerate(prs.slides):
        if len(slides) >= MAX_SLIDES:
            warnings.append(f"deck truncated to {MAX_SLIDES} slides")
            break
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = "\n".join(
                p.text for p in shape.text_frame.paragraphs
            ).strip()
            if not txt:
                continue
            is_title = False
            try:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    is_title = True
            except Exception:
                pass
            if is_title and not title:
                title = txt.splitlines()[0]
            else:
                body_parts.append(txt)
        if not title and body_parts:
            warnings.append(f"slide {i + 1}: no title placeholder")
        slides.append({
            "title": _clip(title, _TITLE_CAP),
            "body": _clip("\n".join(body_parts), _BODY_CAP),
        })
    return slides, warnings


def _extract_pdf_text(file_bytes: bytes):
    try:
        from pypdf import PdfReader
    except ImportError as e:  # pragma: no cover
        raise DeckParseError(f"pdf parser unavailable: {e}")
    try:
        reader = PdfReader(BytesIO(file_bytes))
    except Exception as e:
        raise DeckParseError(f"could not open pdf: {e}")
    slides: list[dict] = []
    warnings: list[str] = []
    for i, page in enumerate(reader.pages):
        if len(slides) >= MAX_SLIDES:
            warnings.append(f"deck truncated to {MAX_SLIDES} slides")
            break
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        if not lines:
            slides.append({"title": "", "body": ""})
            warnings.append(f"slide {i + 1}: no extractable text (image-only?)")
            continue
        slides.append({
            "title": _clip(lines[0], _TITLE_CAP),
            "body": _clip("\n".join(lines[1:]), _BODY_CAP),
        })
    return slides, warnings


# soffice install locations. The Railway runtime PATH often omits the apt/Nix
# dirs (same reason railway-web.sh hunts ffmpeg), and apt's libreoffice-core
# sometimes lacks the /usr/bin/soffice symlink — so probe the real binary at
# /usr/lib/.../program/soffice and the Nix profiles too, not just PATH.
_SOFFICE_CANDIDATES = (
    "/usr/bin/soffice",
    "/usr/bin/libreoffice",
    "/usr/local/bin/soffice",
    "/usr/lib/libreoffice/program/soffice",
    "/opt/libreoffice/program/soffice",
)


def _find_soffice():
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    for p in _SOFFICE_CANDIDATES:
        if os.path.exists(p):
            return p
    # Nix profile dirs (LibreOffice installed via nixPkgs).
    for base in (
        os.path.expanduser("~/.nix-profile/bin"),
        "/root/.nix-profile/bin",
        "/nix/var/nix/profiles/default/bin",
    ):
        cand = os.path.join(base, "soffice")
        if os.path.exists(cand):
            return cand
    return None


def _pptx_to_pdf(file_bytes: bytes) -> bytes:
    soffice = _find_soffice()
    if not soffice:
        raise DeckParseError("LibreOffice not available for PPTX conversion")
    with tempfile.TemporaryDirectory() as td:
        src = os.path.join(td, "in.pptx")
        with open(src, "wb") as f:
            f.write(file_bytes)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf",
                 "--outdir", td, src],
                check=True, capture_output=True, timeout=_SOFFICE_TIMEOUT_SEC,
                # soffice needs a writable HOME for its profile; isolate it per
                # call so concurrent conversions don't fight over the profile.
                env={**os.environ, "HOME": td},
            )
        except subprocess.TimeoutExpired:
            raise DeckParseError("PPTX conversion timed out")
        except subprocess.CalledProcessError as e:
            detail = (e.stderr or b"")[:200].decode("utf-8", "replace")
            raise DeckParseError(f"PPTX conversion failed: {detail}")
        out = os.path.join(td, "in.pdf")
        if not os.path.exists(out):
            raise DeckParseError("PPTX conversion produced no PDF")
        with open(out, "rb") as f:
            return f.read()


def warmup_soffice() -> None:
    """Prime the soffice profile (first conversion is slow — profile init).
    Best-effort; called from gunicorn_conf.py post_worker_init like librosa."""
    try:
        from pptx import Presentation
        buf = BytesIO()
        Presentation().save(buf)
        _pptx_to_pdf(buf.getvalue())
        logger.info("deck_parser: soffice warmup ok")
    except Exception as e:
        logger.warning("deck_parser: soffice warmup skipped: %s", e)
